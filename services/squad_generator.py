from pptx import Presentation
from io import BytesIO


class SquadProposalGenerator:

    def __init__(self, template_path):
        self.prs = Presentation(template_path)

    async def generate(self, data, logo):
        await self._update_logo(logo)
        self._update_client_name(data["cliente"]["nome"])
        self._handle_squad_composition(data["cliente"]["briefing"])

        output_path = "output/proposta_squad.pptx"
        self.prs.save(output_path)
        return output_path

    async def _update_logo(self, logo_file):
        image_bytes = await logo_file.read()
        image_stream = BytesIO(image_bytes)

        slide = self.prs.slides[0]

        for shape in slide.shapes:
            if shape.name == "CLIENT_LOGO":
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                slide.shapes._spTree.remove(shape._element)

                slide.shapes.add_picture(
                    image_stream,
                    left,
                    top,
                    width=width,
                    height=height
                )

    def _update_client_name(self, client_name):
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "<NOME_EMPRESA>" in run.text:
                                run.text = run.text.replace(
                                    "<NOME_EMPRESA>",
                                    client_name
                                )


    def _handle_squad_composition(self, briefing):

        slides_to_remove = []

        for slide in self.prs.slides:
            for shape in slide.shapes:
                title = shape.name if shape.name else ""

                mapping = {
                    "COMPOSICAO_PO": briefing.get("po", "0"),
                    "COMPOSICAO_DEV": briefing.get("dev", "0"),
                    "COMPOSICAO_UX": briefing.get("ux", "0"),
                    "COMPOSICAO_CURADOR": briefing.get("curador", "0"),
                    "COMPOSICAO_ANALISTA": briefing.get("dados", "0"),
                }

                if title in mapping:
                    percentual = mapping[title]

                    if percentual == "0" or percentual is None:
                        slides_to_remove.append(slide)
                    else:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if "<HRS>" in run.text:
                                            run.text = run.text.replace(
                                                "<HRS>",
                                                f"{percentual}H"
                                            )

        for slide in slides_to_remove:
            self._remove_slide(slide)

    def _remove_slide(self, slide):
        slide_id = slide.slide_id
        slides = self.prs.slides._sldIdLst
        for sld in slides:
            if int(sld.get("id")) == slide_id:
                slides.remove(sld)
                break
