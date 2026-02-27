from pptx import Presentation
from io import BytesIO
from enum import Enum
from typing import TypedDict

class PLANS(Enum):
    STARTER_PLAN = "starter"
    SILVER_PLAN = "silver"
    GOLD_PLAN = "gold"
    DIAMOND_PLAN = "diamond"

class AdequatePlanPayload(TypedDict):
    adequatePlan: PLANS

class Client(TypedDict):
    name: str
    briefing: AdequatePlanPayload

class ServiceData(TypedDict):
    tipoProposta: str
    cliente: Client

class SustentationProposalGenerator:

    def __init__(self, template_path: str):
        self.prs = Presentation(template_path)

    async def generate(self, data: ServiceData, logo):
        await self._update_logo(logo)
        self._handle_sustentation_plan(data["cliente"]["briefing"])

        output_path = "output/proposta_sustentacao.pptx"
        self.prs.save(output_path)
        return output_path

    async def _update_logo(self, logo_file):
        image_bytes = await logo_file.read()
        image_stream = BytesIO(image_bytes)

        slide = self.prs.slides[0]

        for shape in slide.shapes:
            if shape.name == "CLIENT_LOGO":
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                slide.shapes._spTree.remove(shape._element)

                slide.shapes.add_picture(
                    image_stream,
                    left,
                    top,
                    width=width,
                    height=height
                )
    
    def _handle_sustentation_plan(self, briefing: AdequatePlanPayload):
        valid_plans = {plan.name for plan in PLANS}

        adequate_plan = briefing.get("adequatePlan")

        if not adequate_plan:
            return

        adequate_plan = adequate_plan.upper() + "_PLAN"

        slides_to_remove = []

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.name in valid_plans and shape.name != adequate_plan:
                    slides_to_remove.append(slide)
                    break

        for slide in slides_to_remove:
            self._remove_slide(slide)

    def _remove_slide(self, slide):
        slide_id = slide.slide_id
        slides = self.prs.slides._sldIdLst
        for sld in slides:
            if int(sld.get("id")) == slide_id:
                slides.remove(sld)
                break

