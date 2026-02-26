from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from enum import Enum
from typing import TypedDict
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import re

class PLANS(Enum):
    STARTER_PLAN = "starter"
    SILVER_PLAN = "silver"
    GOLD_PLAN = "gold"
    DIAMOND_PLAN = "diamond"

class Timeline(TypedDict):
    flowDrawing: str
    drawingHomologation: str
    development: str
    qaHomologation: str
    clientHomologation: str

class AdequatePlanPayload(TypedDict):
    mainGoal: str
    briefingDetails: list[str]
    timeLine: Timeline
    adequatePlan: PLANS

class Client(TypedDict):
    name: str
    briefing: AdequatePlanPayload

class ServiceData(TypedDict):
    tipoProposta: str
    cliente: Client

class GraphService:

    def __init__(self):
        self.prs = Presentation()

    def extrair_semanas(self, texto):
        """Extrai o número de semanas de strings como '2 semanas' ou '1 semana'"""
        if not texto: return 0
        match = re.search(r'(\d+)', str(texto))
        return int(match.group(1)) if match else 0

    def formatar_tabela_transparente(self, table, etapas):
        def set_cell_border(cell, color_str="DCDCDC"):
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            
            for side in ['lnL', 'lnR', 'lnT', 'lnB']:
                # Localiza ou cria o elemento da borda (ex: lnL)
                ln = tcPr.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{side}')
                if ln is None:
                    ln = OxmlElement(f'a:{side}')
                    tcPr.append(ln)
                
                # Define a largura (w)
                ln.set('w', '6350') # Aprox 0.5pt
                ln.set('cap', 'flat')
                ln.set('cmpd', 'sng')
                ln.set('algn', 'ctr')

                # LIMPA filhos antigos para evitar conflitos
                for child in ln.getchildren():
                    ln.remove(child)

                # 1. DEFINE A COR
                srgbClr = OxmlElement('a:srgbClr')
                srgbClr.set('val', color_str)
                ln.append(srgbClr)

                # 2. O PULO DO GATO: DEFINE O ESTILO DA LINHA COMO SÓLIDO
                # Sem isso, o PPT pode entender que a linha existe mas é "invisível"
                prstDash = OxmlElement('a:prstDash')
                prstDash.set('val', 'solid')
                ln.append(prstDash)

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell.fill.background()
                set_cell_border(cell, "DCDCDC") # Cinza claro

                # Manter a lógica de texto que já funciona
                if col_idx == 0 and row_idx > 0:
                    cor_etapa = etapas[row_idx-1][2]
                    p = cell.text_frame.paragraphs[0]
                    p.font.name = 'Roboto'
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = cor_etapa
    
    def adicionar_barra_com_padding(self, slide, table_shape, row_idx, col_idx, cor_rgb):
        """
        Calcula a posição usando o shape da tabela (GraphicFrame) para evitar erros de atributo.
        """
        table = table_shape.table
        
        # 1. Posição inicial baseada no Shape que contém a tabela
        left_pos = table_shape.left 
        top_pos = table_shape.top
            
        # 2. Somar larguras das colunas anteriores
        for c in range(col_idx):
            left_pos += table.columns[c].width
            
        # 3. Somar alturas das linhas anteriores
        for r in range(row_idx):
            top_pos += table.rows[r].height
            
        # 4. Dimensões da célula alvo
        cell_width = table.columns[col_idx].width
        cell_height = table.rows[row_idx].height
        
        # 5. Padding (Ajuste conforme o gosto visual)
        padding_h = Pt(8) # Padding horizontal
        padding_v = Pt(6) # Padding vertical
        
        bar_left = left_pos + padding_h
        bar_top = top_pos + padding_v
        bar_width = cell_width - (padding_h * 2)
        bar_height = cell_height - (padding_v * 2)
        
        # 6. Adicionar o retângulo arredondado
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top, bar_width, bar_height
        )
        
        # Estilização final
        rect.fill.solid()
        rect.fill.fore_color.rgb = cor_rgb
        rect.line.fill.background() # Remove contorno
        rect.adjustments[0] = 0.2    # Suaviza o arredondamento

    def generate(self, data: ServiceData):
        timeline_data = data["cliente"]["briefing"]["timeLine"]
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = "Timeline do Projeto"

        etapas = [
            ("flowDrawing", "DESENHO DE FLUXO", RGBColor(31, 73, 125)),
            ("drawingHomologation", "HOMOLOGAÇÃO DO DESENHO", RGBColor(255, 51, 153)),
            ("development", "DESENVOLVIMENTO DO FLUXO", RGBColor(128, 110, 180)),
            ("qaHomologation", "TESTES INTERNOS (QAs)", RGBColor(0, 0, 255)),
            ("clientHomologation", "HOMOLOGAÇÃO CLIENTE", RGBColor(0, 176, 240))
        ]

        total_semanas = sum(self.extrair_semanas(timeline_data.get(k, 0)) for k, _, _ in etapas)
        num_colunas = max(total_semanas, 6) + 1 

        rows, cols = len(etapas) + 1, num_colunas
        left, top = Inches(0.5), Inches(2.0)
        width, height = Inches(9.0), Inches(3.5)
        
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        table.style = "TableGrid"
        
        # Limpar fundo e bordas antes de escrever
        self.formatar_tabela_transparente(table, etapas)
        
        # 3. Re-inserir Cabeçalhos (SEM 1, SEM 2...)
        for i in range(1, cols):
            cell = table.cell(0, i)
            cell.text = f"SEM {i}"
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            # Forçar visibilidade e fonte
            run = p.runs[0]
            run.font.name = 'Roboto'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(120, 120, 120) # Cinza discreto para os cabeçalhos

        # 4. Loop de Preenchimento das Etapas
        coluna_atual = 1
        for row_idx, (chave_json, nome_etapa, cor) in enumerate(etapas, start=1):
            cell_nome = table.cell(row_idx, 0)
            cell_nome.text = nome_etapa
            
            # Formatação Roboto para as nomenclaturas laterais
            p = cell_nome.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = 'Roboto'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = cor # Cor correspondente à etapa

            duracao = self.extrair_semanas(timeline_data.get(chave_json, 0))
            
            for _ in range(duracao):
                if coluna_atual < cols:
                    self.adicionar_barra_com_padding(slide, shape, row_idx, coluna_atual, cor)
                    coluna_atual += 1

        output_path = "output/proposta_teste.pptx"
        self.prs.save(output_path)
        return output_path