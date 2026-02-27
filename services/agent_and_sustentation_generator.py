from pptx import Presentation
from io import BytesIO
from enum import Enum
from typing import TypedDict
from copy import deepcopy
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import logging
import io
import math

logger = logging.getLogger(__name__)

class PLANS(Enum):
    STARTER_PLAN = "starter"
    SILVER_PLAN = "silver"
    GOLD_PLAN = "gold"
    DIAMOND_PLAN = "diamond"

class Timeline(TypedDict):
    flowDrawing: float
    drawingHomologation: str
    development: str
    qaHomologation: str
    clientHomologation: str

class AdequatePlanPayload(TypedDict):
    mainGoal: str
    briefingDetails: list[str]
    timeLine: Timeline
    adequatePlan: PLANS

class Client(TypedDict):
    name: str
    briefing: AdequatePlanPayload

class ServiceData(TypedDict):
    tipoProposta: str
    cliente: Client


class AgentAndSustentationProposalGenerator:

    def __init__(self, template_path: str):
        self.prs = Presentation(template_path)

    async def generate(self, data: ServiceData, logo):
        await self._update_logo(logo)
        self._handle_project_scope(data["cliente"]["briefing"])
        self._handle_project_timeline(data["cliente"]["briefing"])
        self._handle_sustentation_plan(data["cliente"]["briefing"])

        output_path = "output/proposta_agent_sustentacao.pptx"
        self.prs.save(output_path)
        return output_path

    async def _update_logo(self, logo_file):
        logger.info(f"Iniciando a atualização da logo...")


        image_bytes = await logo_file.read()
        image_stream = BytesIO(image_bytes)

        slide = self.prs.slides[0]

        for shape in slide.shapes:
            if shape.name == "CLIENT_LOGO":
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                slide.shapes._spTree.remove(shape._element)

                slide.shapes.add_picture(
                    image_stream,
                    left,
                    top,
                    width=width,
                    height=height
                )
    
    def _normalized_briefing_details(self, detail: str) -> str:
        return f"- {detail.strip()}"
    
    def _chunk_briefing(self, details: list[str], limit: int = 500):
        chunks = []
        current_chunk = ""
        
        for detail in details:
            formatted = f"{detail.strip()}\n"
            
            if len(current_chunk) + len(formatted) > limit:
                chunks.append(current_chunk.strip())
                current_chunk = formatted
            else:
                current_chunk += formatted

        if current_chunk:
            chunks.append(current_chunk.strip())

        return chunks

    def _duplicate_slide(self, slide):
        slide_layout = next(
            layout for layout in self.prs.slide_layouts
            if layout.name.lower() == "blank"
        )

        new_slide = self.prs.slides.add_slide(slide_layout)

        # remover placeholders herdados
        for shape in list(new_slide.shapes):
            if shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)

        IMAGE_SHAPES = {"PINK_IMAGE", "DIGITALBOT_LOGO"}
        RECTANGLE_SHAPES = {"SCOPE", "SCOPE_MAIN_GOAL", "SCOPE_DETAILS"}

        for shape in slide.shapes:
            # 🔹 Imagens
            if shape.name in IMAGE_SHAPES:
                image_stream = io.BytesIO(shape.image.blob)

                new_picture = new_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )

                new_picture.name = shape.name

            # 🔹 Retângulos vazios
            elif shape.name in RECTANGLE_SHAPES:
                # Clonar o shape original inteiro com todas as propriedades
                new_element = deepcopy(shape._element)
                new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
                
                if shape.name != "SCOPE":
                    # Encontrar o shape adicionado e limpar seu texto
                    for new_shape in new_slide.shapes:
                        if new_shape.name == shape.name:
                            new_shape.text_frame.clear()
                            break

        # 🔹 Reordenar slide
        slide_id_list = self.prs.slides._sldIdLst

        for idx, sldId in enumerate(slide_id_list):
            if sldId.id == slide.slide_id:
                original_index = idx
                break

        new_sldId = slide_id_list[-1]
        slide_id_list.remove(new_sldId)
        slide_id_list.insert(original_index + 1, new_sldId)

        return new_slide

    def _handle_project_scope(self, briefing):
        logger.info(f"Iniciando a atualização dos slides de escopo...")

        main_goal = briefing.get("mainGoal")
        briefing_details = briefing.get("briefingDetails")

        if not briefing_details:
            return

        chunks = self._chunk_briefing(briefing_details, 500)

        scope_slide = None

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.name == "SCOPE":
                    scope_slide = slide
                    break
            if scope_slide:
                break
        
        if not scope_slide:
            return

        slides_to_fill = [scope_slide]

        for _ in range(len(chunks) - 1):
            duplicated = self._duplicate_slide(scope_slide)
            slides_to_fill.append(duplicated)

        for slide, chunk in zip(slides_to_fill, chunks):

            for shape in slide.shapes:
                if shape.name == "SCOPE_MAIN_GOAL":
                    text_frame = shape.text_frame
                    text_frame.clear()

                    p = text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT
                    run = p.add_run()
                    run.text = f"🎯 Objetivo Geral: {main_goal}"
                    run.font.name = "Lexend"
                    run.font.size = Pt(22)
                    run.font.color.rgb = RGBColor(0, 0, 0)

                if shape.name == "SCOPE_DETAILS":
                    text_frame = shape.text_frame
                    text_frame.clear()

                    lines = [l.strip().lstrip("-").strip() 
                     for l in chunk.split("\n") if l.strip()]

                    if not lines:
                        continue
                    
                    first_p = text_frame.paragraphs[0]
                    first_p.text = f"• {lines[0]}"
                    first_p.level = 0
                    first_p.alignment = PP_ALIGN.LEFT
                    first_p.runs[0].font.name = "Lexend"
                    first_p.runs[0].font.size = Pt(18)
                    first_p.runs[0].font.color.rgb = RGBColor(0, 0, 0)

                    space_p = text_frame.add_paragraph()
                    space_p.text = ""

                    for line in lines[1:]:
                        p = text_frame.add_paragraph()
                        p.text = f"• {line}"
                        p.level = 0
                        p.alignment = PP_ALIGN.LEFT
                        p.runs[0].font.name = "Lexend"
                        p.runs[0].font.size = Pt(18)
                        p.runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        
                        # Parágrafo vazio entre linhas
                        space_p = text_frame.add_paragraph()
                        space_p.text = ""

    def _remove_old_bars(self, slide):
        shapes_to_remove = []

        for shape in slide.shapes:
            if shape.name.startswith("BAR_"):
                shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def _handle_project_timeline(self, briefing: AdequatePlanPayload):
        logger.info("Construindo timeline estilo Gantt...")

        timeline_data = briefing["timeLine"]

        etapas = [
            ("flowDrawing", RGBColor(31, 73, 125)),
            ("drawingHomologation", RGBColor(255, 51, 153)),
            ("development", RGBColor(128, 110, 180)),
            ("qaHomologation", RGBColor(0, 0, 255)),
            ("clientHomologation", RGBColor(0, 176, 240)),
        ]

        # localizar tabela
        timeline_slide = None
        table_shape = None

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.name == "GRAPH_SHAPE":
                    timeline_slide = slide
                    table_shape = shape
                    break
            if timeline_slide:
                break

        if not timeline_slide:
            logger.warning("GRAPH_SHAPE não encontrada.")
            return

        table = table_shape.table

        # remover barras antigas
        self._remove_old_bars(timeline_slide)

        # total semanas (considerando fração)
        total_semanas = sum(
            float(timeline_data.get(k, 0) or 0) for k, _ in etapas
        )

        total_semanas = max(math.ceil(total_semanas), 6)

        # cabeçalho
        for col in range(1, total_semanas + 1):
            cell = table.cell(0, col)
            cell.text = f"SEM {col}"

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run = p.runs[0]
            run.font.name = "Roboto"
            run.font.size = Pt(16)
            run.font.italic = True
            run.font.color.rgb = RGBColor(120, 120, 120)

        coluna_atual = 1.0  # agora é float

        for row_idx, (key, cor) in enumerate(etapas, start=1):

            duracao = float(timeline_data.get(key, 0) or 0)

            if duracao <= 0:
                continue

            # posição base
            left = table_shape.left
            top = table_shape.top

            # deslocamento horizontal (parte inteira)
            col_inteira = int(coluna_atual)

            for c in range(col_inteira):
                left += table.columns[c].width

            # parte fracionada inicial
            fracao_inicio = coluna_atual - col_inteira
            if fracao_inicio > 0:
                left += table.columns[col_inteira].width * fracao_inicio

            # deslocamento vertical
            for r in range(row_idx):
                top += table.rows[r].height

            # largura considerando fração
            largura_total = 0

            parte_inteira = int(duracao)
            fracao = duracao - parte_inteira

            # somar colunas inteiras
            for i in range(parte_inteira):
                largura_total += table.columns[col_inteira + i].width

            # parte fracionada final
            if fracao > 0:
                largura_total += table.columns[col_inteira + parte_inteira].width * fracao

            altura = table.rows[row_idx].height

            padding_h = Pt(6)
            padding_v = Pt(4)

            barra = timeline_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left + padding_h,
                top + padding_v,
                largura_total - (padding_h * 2),
                altura - (padding_v * 2),
            )

            barra.name = f"BAR_{key}"

            barra.fill.solid()
            barra.fill.fore_color.rgb = cor
            barra.line.fill.background()
            barra.adjustments[0] = 0.3

            coluna_atual += duracao
                

    def _handle_sustentation_plan(self, briefing: AdequatePlanPayload):
        logger.info(f"Iniciando a escolha de slide de plano de sustentação...")

        valid_plans = {plan.name for plan in PLANS}

        adequate_plan = briefing.get("adequatePlan")

        if not adequate_plan:
            return

        adequate_plan = adequate_plan.upper() + "_PLAN"

        slides_to_remove = []

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.name in valid_plans and shape.name != adequate_plan:
                    slides_to_remove.append(slide)
                    break

        for slide in slides_to_remove:
            self._remove_slide(slide)

    def _remove_slide(self, slide):
        slide_id = slide.slide_id
        slides = self.prs.slides._sldIdLst
        for sld in slides:
            if int(sld.get("id")) == slide_id:
                slides.remove(sld)
                break

